#!/usr/bin/env python3
"""
统一素材读取器 —— teacher-paper skill 自包含组件
支持：.docx .pdf .pptx .xlsx .xls .csv .txt .md  及图片提示
用法：
    python3 read_material.py <文件路径> [文件路径2 ...]
输出：纯文本（表格转为 " | " 分隔），多文件以分隔线拼接，打到 stdout。

设计原则：
- 缺少依赖库时给出可执行的 pip 安装提示，而非崩溃。
- 任何电脑均可运行，不依赖其他 skill 或硬编码路径。
"""
import sys
import os

IMG_EXT = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp", ".tiff"}


def _need(pkg):
    return (f"[缺少依赖] 读取此文件需要 Python 库 `{pkg}`。\n"
            f"请运行：pip3 install {pkg}\n"
            f"（或让助手改用其它方式读取。）")


def read_docx(path):
    try:
        import docx
    except ImportError:
        return _need("python-docx")
    from docx.document import Document as _Doc
    from docx.oxml.table import CT_Tbl
    from docx.oxml.text.paragraph import CT_P
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    d = docx.Document(path)
    out = []
    body = d.element.body
    for child in body.iterchildren():
        if isinstance(child, CT_P):
            t = Paragraph(child, d).text.strip()
            if t:
                out.append(t)
        elif isinstance(child, CT_Tbl):
            for row in Table(child, d).rows:
                cells = [c.text.strip() for c in row.cells]
                if any(cells):
                    out.append(" | ".join(cells))
    note = ""
    # 提示可能存在的图片（阅读材料常以图片嵌入）
    try:
        import zipfile
        z = zipfile.ZipFile(path)
        media = [n for n in z.namelist() if n.startswith("word/media/")]
        if media:
            note = (f"\n\n[提示] 该 docx 内含 {len(media)} 张嵌入图片"
                    f"（{', '.join(os.path.basename(m) for m in media[:8])}…）。"
                    f"如阅读材料/古诗文以图片形式存在，文本可能不全，"
                    f"请解压 word/media/ 后用多模态读取图片。")
    except Exception:
        pass
    return "\n".join(out) + note


def read_pdf(path):
    try:
        import pdfplumber
    except ImportError:
        try:
            from pypdf import PdfReader
            r = PdfReader(path)
            return "\n".join((p.extract_text() or "") for p in r.pages)
        except ImportError:
            return _need("pdfplumber") + "\n（或 pip3 install pypdf）"
    out = []
    with pdfplumber.open(path) as pdf:
        for i, page in enumerate(pdf.pages, 1):
            out.append(f"\n--- 第{i}页 ---")
            txt = page.extract_text() or ""
            out.append(txt)
            for tbl in page.extract_tables():
                for row in tbl:
                    out.append(" | ".join((c or "") for c in row))
    return "\n".join(out)


def read_pptx(path):
    try:
        from pptx import Presentation
    except ImportError:
        return _need("python-pptx")
    prs = Presentation(path)
    out = []
    for i, slide in enumerate(prs.slides, 1):
        out.append(f"\n--- 幻灯片{i} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = "".join(run.text for run in para.runs).strip()
                    if t:
                        out.append(t)
            if shape.has_table:
                for row in shape.table.rows:
                    out.append(" | ".join(c.text.strip() for c in row.cells))
    return "\n".join(out)


def read_xlsx(path):
    try:
        from openpyxl import load_workbook
    except ImportError:
        return _need("openpyxl")
    wb = load_workbook(path, data_only=True)
    out = []
    for ws in wb.worksheets:
        out.append(f"\n--- 工作表：{ws.title} ---")
        for row in ws.iter_rows(values_only=True):
            cells = ["" if v is None else str(v) for v in row]
            if any(cells):
                out.append(" | ".join(cells))
    return "\n".join(out)


def read_csv(path):
    """Bug-B4 修复：先用 chardet/charset_normalizer 探测编码（utf-8 字节用 gbk
    解码常不抛异常而是产生乱码），避免按 utf-8/gbk 静默匹配错误编码。
    探测不到时回退原有 utf-8-sig/gbk/utf-8 顺序，但增加'非可打印率'校验。"""
    import csv
    # 优先用 chardet 或 charset_normalizer 探测
    enc = None
    try:
        with open(path, "rb") as f:
            sample = f.read(8192)
        try:
            import chardet
            res = chardet.detect(sample)
            if res and res.get("confidence", 0) > 0.7:
                enc = res.get("encoding")
        except ImportError:
            try:
                from charset_normalizer import from_bytes
                result = from_bytes(sample).best()
                if result:
                    enc = result.encoding
            except ImportError:
                pass
    except OSError:
        pass

    encodings = [enc] if enc else []
    encodings += ["utf-8-sig", "gbk", "utf-8"]

    def _gibberish_ratio(text):
        """U+FFFD / 控制字符占比；高于 1% 视为乱码。"""
        if not text:
            return 0
        bad = sum(1 for c in text if c == "�" or (ord(c) < 32 and c not in "\r\n\t"))
        return bad / len(text)

    for e in encodings:
        if not e:
            continue
        try:
            with open(path, newline="", encoding=e, errors="replace") as f:
                text = f.read()
            if _gibberish_ratio(text) < 0.01:
                out = []
                with open(path, newline="", encoding=e, errors="replace") as f:
                    for row in csv.reader(f):
                        out.append(" | ".join(row))
                return "\n".join(out)
        except (UnicodeDecodeError, OSError):
            continue
    return "[无法识别 CSV 编码，请改存为 UTF-8 后重试，或安装 chardet：pip3 install chardet]"


def read_text(path):
    for enc in ("utf-8", "gbk", "utf-16"):
        try:
            with open(path, encoding=enc) as f:
                return f.read()
        except UnicodeDecodeError:
            continue
    return "[无法识别文本编码]"


def read_one(path):
    if not os.path.exists(path):
        return f"[文件不存在] {path}"
    ext = os.path.splitext(path)[1].lower()
    if ext == ".docx":
        return read_docx(path)
    if ext == ".pdf":
        return read_pdf(path)
    if ext == ".pptx":
        return read_pptx(path)
    if ext in (".xlsx", ".xlsm"):
        return read_xlsx(path)
    if ext == ".csv":
        return read_csv(path)
    if ext in (".txt", ".md", ".markdown"):
        return read_text(path)
    if ext in IMG_EXT:
        return (f"[图片文件] {path}\n"
                f"图片需用助手的多模态 Read 工具直接识别，本脚本不处理。")
    if ext in (".doc", ".ppt", ".xls"):
        return (f"[旧版二进制 Office 格式] {path}\n"
                f"请先另存为 .docx/.pptx/.xlsx，或安装 libreoffice 转换：\n"
                f"  soffice --headless --convert-to docx \"{path}\"")
    return read_text(path)  # 兜底按文本读


def main():
    if len(sys.argv) < 2:
        print(__doc__)
        sys.exit(1)
    parts = []
    for p in sys.argv[1:]:
        parts.append(f"\n{'='*60}\n# 来源：{os.path.basename(p)}\n{'='*60}")
        parts.append(read_one(p))
    print("\n".join(parts))


if __name__ == "__main__":
    main()
